from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from loguru import logger

from .loader_registry import LoaderRegistry, DocumentContent


@LoaderRegistry.register([".pptx"])
class PptxLoader:
    def load(self, file_path: Path) -> DocumentContent:
        prs = Presentation(str(file_path))
        parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(cells))
            if slide_text:
                parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

        full_text = "\n\n".join(parts)
        logger.info(f"PPTX loaded: {file_path.name}, {len(full_text)} chars")
        return DocumentContent(text=full_text)
